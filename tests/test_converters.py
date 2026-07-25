"""Coverage for the extensible common-format converter registry."""

from __future__ import annotations

import json
import zipfile

from pdf2mdx.converters import convert_file, supported_extensions


def test_registry_contains_primary_formats() -> None:
    expected = {".pdf", ".xmind", ".docx", ".pptx", ".xlsx", ".html", ".txt", ".csv"}
    assert expected.issubset(set(supported_extensions()))


def test_text_and_csv_conversion(tmp_path) -> None:
    text = tmp_path / "笔记.txt"
    text.write_text("第一行\n第二行", encoding="utf-8")
    assert "第一行" in convert_file(str(text), str(tmp_path))

    csv_file = tmp_path / "数据.csv"
    csv_file.write_text("姓名,数量\n墨桥,3", encoding="utf-8-sig")
    markdown = convert_file(str(csv_file), str(tmp_path))
    assert "| 姓名 | 数量 |" in markdown
    assert "| 墨桥 | 3 |" in markdown


def test_json_and_xml_conversion(tmp_path) -> None:
    json_file = tmp_path / "data.json"
    json_file.write_text(json.dumps({"name": "墨桥"}, ensure_ascii=False), encoding="utf-8")
    assert '"name": "墨桥"' in convert_file(str(json_file), str(tmp_path))

    xml_file = tmp_path / "data.xml"
    xml_file.write_text("<root><item id='1'>内容</item></root>", encoding="utf-8")
    markdown = convert_file(str(xml_file), str(tmp_path))
    assert "**root**" in markdown
    assert "**item**" in markdown
    assert "内容" in markdown


def test_docx_xlsx_and_pptx_conversion(tmp_path) -> None:
    from docx import Document
    from openpyxl import Workbook
    from pptx import Presentation

    doc_path = tmp_path / "sample.docx"
    doc = Document()
    doc.add_heading("章节", level=1)
    doc.add_paragraph("正文")
    table = doc.add_table(rows=2, cols=2)
    table.cell(0, 0).text = "列A"
    table.cell(0, 1).text = "列B"
    table.cell(1, 0).text = "甲"
    table.cell(1, 1).text = "乙"
    doc.save(doc_path)
    doc_md = convert_file(str(doc_path), str(tmp_path))
    assert "## 章节" in doc_md and "| 列A | 列B |" in doc_md

    sheet_path = tmp_path / "sample.xlsx"
    book = Workbook()
    sheet = book.active
    sheet.title = "明细"
    sheet.append(["名称", "数量"])
    sheet.append(["墨桥", 2])
    book.save(sheet_path)
    xlsx_md = convert_file(str(sheet_path), str(tmp_path))
    assert "## 明细" in xlsx_md and "| 墨桥 | 2 |" in xlsx_md

    deck_path = tmp_path / "sample.pptx"
    deck = Presentation()
    slide = deck.slides.add_slide(deck.slide_layouts[5])
    slide.shapes.title.text = "演示标题"
    deck.save(deck_path)
    pptx_md = convert_file(str(deck_path), str(tmp_path))
    assert "## 第 1 页" in pptx_md and "演示标题" in pptx_md


def test_html_conversion_preserves_table(tmp_path) -> None:
    source = tmp_path / "page.html"
    source.write_text(
        "<html><head><title>页面</title></head><body><main><div>"
        "<h2>说明</h2><p>访问 <a href='https://example.com'>示例</a></p>"
        "<table><tr><th>A</th><th>B</th></tr>"
        "<tr><td>1</td><td>2</td></tr></table></div></main></body></html>",
        encoding="utf-8",
    )
    markdown = convert_file(str(source), str(tmp_path))
    assert "# 页面" in markdown
    assert "### 说明" in markdown
    assert "[示例](https://example.com)" in markdown
    assert "| A | B |" in markdown


def test_xmind_embedded_image_uses_asset_directory(tmp_path) -> None:
    source = tmp_path / "map.xmind"
    content = [{
        "title": "画布",
        "rootTopic": {
            "title": "根主题",
            "image": {"src": "xap:resources/picture.png"},
            "children": {"attached": []},
        },
    }]
    with zipfile.ZipFile(source, "w") as archive:
        archive.writestr("content.json", json.dumps(content, ensure_ascii=False))
        archive.writestr("resources/picture.png", b"\x89PNG\r\n\x1a\n")
    markdown = convert_file(str(source), str(tmp_path))
    assert "map_assets/picture.png" in markdown
    assert (tmp_path / "map_assets" / "picture.png").exists()
