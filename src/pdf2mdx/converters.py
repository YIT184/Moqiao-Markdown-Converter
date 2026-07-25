"""Extensible converter registry and converters for common office/text formats."""

from __future__ import annotations

import csv
import html
import json
import os
import pathlib
import shutil
import zipfile
from collections.abc import Callable
from typing import Any
from urllib.parse import unquote, urlparse
import xml.etree.ElementTree as ET

from .markdown_utils import collision_free_name, normalize_table, safe_stem
from .pdf_converter import PdfConverter
from .xmind_converter import XmindConverter

Converter = Callable[[str, str, dict[str, Any]], str]
REGISTRY: dict[str, Converter] = {}


def register(*extensions: str) -> Callable[[Converter], Converter]:
    """Register a conversion function for one or more extensions."""
    def decorator(func: Converter) -> Converter:
        for extension in extensions:
            REGISTRY[extension.lower()] = func
        return func
    return decorator


def supported_extensions() -> tuple[str, ...]:
    return tuple(sorted(REGISTRY))


def convert_file(path: str, output_dir: str, options: dict[str, Any] | None = None) -> str:
    """Convert a supported file and return Markdown without writing the .md file."""
    extension = pathlib.Path(path).suffix.lower()
    converter = REGISTRY.get(extension)
    if converter is None:
        raise ValueError(f"暂不支持 {extension or '无扩展名'} 文件")
    return converter(path, output_dir, options or {})


def _asset_dir(path: str, output_dir: str) -> pathlib.Path:
    target = pathlib.Path(output_dir) / f"{safe_stem(pathlib.Path(path).stem)}_assets"
    target.mkdir(parents=True, exist_ok=True)
    return target


def _asset_link(asset_dir: pathlib.Path, filename: str, alt: str = "图片") -> str:
    relative = f"{asset_dir.name}/{filename}".replace("\\", "/")
    return f"![{alt}]({relative})"


@register(".pdf")
def convert_pdf(path: str, output_dir: str, options: dict[str, Any]) -> str:
    converter = PdfConverter(
        preserve_images=bool(options.get("images", True)),
        detect_tables=bool(options.get("tables", True)),
        preserve_vectors=bool(options.get("vectors", True)),
        enhance_line_art=bool(options.get("enhance_line_art", True)),
        password=options.get("password") or None,
    )
    return converter.convert(path, output_dir)


@register(".xmind")
def convert_xmind(path: str, output_dir: str, options: dict[str, Any]) -> str:
    return XmindConverter().convert(path, output_dir)


@register(".txt")
def convert_text(path: str, output_dir: str, options: dict[str, Any]) -> str:
    raw = pathlib.Path(path).read_bytes()
    for encoding in ("utf-8-sig", "utf-16", "gb18030"):
        try:
            text = raw.decode(encoding)
            break
        except UnicodeDecodeError:
            continue
    else:
        text = raw.decode("utf-8", errors="replace")
    return f"# {safe_stem(pathlib.Path(path).stem)}\n\n{text.strip()}\n"


@register(".csv")
def convert_csv(path: str, output_dir: str, options: dict[str, Any]) -> str:
    raw = pathlib.Path(path).read_bytes()
    text = None
    for encoding in ("utf-8-sig", "gb18030", "utf-16"):
        try:
            text = raw.decode(encoding)
            break
        except UnicodeDecodeError:
            continue
    text = text if text is not None else raw.decode("utf-8", errors="replace")
    try:
        dialect = csv.Sniffer().sniff(text[:4096], delimiters=",;\t|")
    except csv.Error:
        dialect = csv.excel
    rows = [[str(cell) for cell in row] for row in csv.reader(text.splitlines(), dialect)]
    return f"# {safe_stem(pathlib.Path(path).stem)}\n\n{normalize_table(rows)}"


@register(".json")
def convert_json(path: str, output_dir: str, options: dict[str, Any]) -> str:
    data = json.loads(pathlib.Path(path).read_text(encoding="utf-8-sig"))
    title = safe_stem(pathlib.Path(path).stem)
    return f"# {title}\n\n```json\n{json.dumps(data, ensure_ascii=False, indent=2)}\n```\n"


@register(".xml")
def convert_xml(path: str, output_dir: str, options: dict[str, Any]) -> str:
    root = ET.parse(path).getroot()
    title = safe_stem(pathlib.Path(path).stem)
    lines = [f"# {title}", ""]

    def walk(node: ET.Element, depth: int) -> None:
        tag = node.tag.split("}", 1)[-1]
        indent = "  " * depth
        value = (node.text or "").strip()
        attributes = " ".join(f'{key}="{value}"' for key, value in node.attrib.items())
        detail = " - ".join(part for part in (attributes, value) if part)
        lines.append(f"{indent}- **{tag}**" + (f": {detail}" if detail else ""))
        for child in node:
            walk(child, depth + 1)

    walk(root, 0)
    return "\n".join(lines) + "\n"


@register(".docx")
def convert_docx(path: str, output_dir: str, options: dict[str, Any]) -> str:
    from docx import Document
    from docx.table import Table
    from docx.text.paragraph import Paragraph

    document = Document(path)
    title = safe_stem(pathlib.Path(path).stem)
    parts = [f"# {title}"]
    for block in document.iter_inner_content():
        if isinstance(block, Paragraph):
            text = block.text.strip()
            if not text:
                continue
            style = (block.style.name if block.style else "").lower()
            if style.startswith("heading"):
                digits = "".join(char for char in style if char.isdigit())
                level = min(6, (int(digits) if digits else 1) + 1)
                parts.append(f"{'#' * level} {text}")
            elif "list" in style:
                parts.append(f"- {text}")
            else:
                parts.append(text)
        elif isinstance(block, Table):
            rows = [[cell.text for cell in row.cells] for row in block.rows]
            table = normalize_table(rows)
            if table:
                parts.append(table.rstrip())

    if options.get("images", True):
        assets = _asset_dir(path, output_dir)
        extracted: list[str] = []
        with zipfile.ZipFile(path) as archive:
            for member in archive.namelist():
                if not member.startswith("word/media/"):
                    continue
                preferred = pathlib.Path(member).name
                name = collision_free_name(str(assets), preferred)
                (assets / name).write_bytes(archive.read(member))
                extracted.append(_asset_link(assets, name, pathlib.Path(name).stem))
        if extracted:
            parts.extend(["## 文档图片", *extracted])
    return "\n\n".join(parts) + "\n"


@register(".pptx")
def convert_pptx(path: str, output_dir: str, options: dict[str, Any]) -> str:
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    deck = Presentation(path)
    title = safe_stem(pathlib.Path(path).stem)
    assets = _asset_dir(path, output_dir)
    parts = [f"# {title}"]
    image_index = 0
    for number, slide in enumerate(deck.slides, 1):
        parts.append(f"## 第 {number} 页")
        for shape in sorted(slide.shapes, key=lambda item: (item.top, item.left)):
            if getattr(shape, "has_table", False):
                rows = [[cell.text for cell in row.cells] for row in shape.table.rows]
                table = normalize_table(rows)
                if table:
                    parts.append(table.rstrip())
            elif shape.shape_type == MSO_SHAPE_TYPE.PICTURE and options.get("images", True):
                image_index += 1
                image = shape.image
                preferred = f"slide_{number:03d}_image_{image_index:03d}.{image.ext}"
                name = collision_free_name(str(assets), preferred)
                (assets / name).write_bytes(image.blob)
                parts.append(_asset_link(assets, name, f"第 {number} 页图片"))
            elif getattr(shape, "has_text_frame", False):
                text = shape.text.strip()
                if text:
                    parts.append(text)
    return "\n\n".join(parts) + "\n"


@register(".xlsx")
def convert_xlsx(path: str, output_dir: str, options: dict[str, Any]) -> str:
    from openpyxl import load_workbook

    workbook = load_workbook(path, data_only=False, read_only=True)
    parts = [f"# {safe_stem(pathlib.Path(path).stem)}"]
    for sheet in workbook.worksheets:
        parts.append(f"## {sheet.title}")
        rows = [
            ["" if cell is None else str(cell) for cell in row]
            for row in sheet.iter_rows(values_only=True)
        ]
        while rows and all(not cell.strip() for cell in rows[-1]):
            rows.pop()
        if rows:
            width = max(len(row) for row in rows)
            while width and all((row[width - 1] if width <= len(row) else "").strip() == "" for row in rows):
                width -= 1
            table = normalize_table([row[:width] for row in rows])
            if table:
                parts.append(table.rstrip())
    workbook.close()
    return "\n\n".join(parts) + "\n"


@register(".html", ".htm")
def convert_html(path: str, output_dir: str, options: dict[str, Any]) -> str:
    from bs4 import BeautifulSoup, NavigableString, Tag

    source = pathlib.Path(path)
    soup = BeautifulSoup(source.read_text(encoding="utf-8", errors="replace"), "html.parser")
    title = soup.title.get_text(" ", strip=True) if soup.title else safe_stem(source.stem)
    parts = [f"# {title}"]
    assets = _asset_dir(path, output_dir)
    def inline_text(node: Tag) -> str:
        chunks: list[str] = []
        for child in node.children:
            if isinstance(child, NavigableString):
                chunks.append(str(child))
            elif isinstance(child, Tag) and child.name.lower() == "a":
                label = child.get_text(" ", strip=True)
                href = child.get("href", "")
                chunks.append(f"[{label}]({href})" if href else label)
            elif isinstance(child, Tag) and child.name.lower() == "br":
                chunks.append("\n")
            elif isinstance(child, Tag):
                chunks.append(inline_text(child))
        return " ".join("".join(chunks).split())

    def walk(node: Any) -> None:
        if isinstance(node, NavigableString):
            text = str(node).strip()
            if text:
                parts.append(text)
            return
        if not isinstance(node, Tag):
            return
        name = node.name.lower()
        text = inline_text(node)
        if name in {"script", "style", "noscript", "title"}:
            return
        if name in {"h1", "h2", "h3", "h4", "h5", "h6"}:
            level = min(6, int(name[1]) + 1)
            parts.append(f"{'#' * level} {text}")
        elif name in {"ul", "ol"}:
            ordered = name == "ol"
            for index, item in enumerate(node.find_all("li", recursive=False), 1):
                parts.append(f"{index}. {item.get_text(' ', strip=True)}" if ordered else f"- {item.get_text(' ', strip=True)}")
        elif name == "table":
            rows = [
                [cell.get_text(" ", strip=True) for cell in row.find_all(["th", "td"])]
                for row in node.find_all("tr")
            ]
            table = normalize_table(rows)
            if table:
                parts.append(table.rstrip())
        elif name == "img" and options.get("images", True):
            src = unquote(urlparse(node.get("src", "")).path)
            if os.name == "nt" and src.startswith("/") and len(src) > 2 and src[2] == ":":
                src = src[1:]
            local = pathlib.Path(src)
            if not local.is_absolute():
                local = source.parent / local
            if local.is_file():
                filename = collision_free_name(str(assets), local.name)
                shutil.copy2(local, assets / filename)
                parts.append(_asset_link(assets, filename, node.get("alt") or "图片"))
        elif name in {"p", "blockquote", "pre"} and text:
            if name == "blockquote":
                parts.append("\n".join(f"> {line}" for line in text.splitlines()))
            elif name == "pre":
                parts.append(f"```\n{node.get_text()}\n```")
            else:
                parts.append(text)
        else:
            for child in node.children:
                walk(child)

    for child in (soup.body or soup).children:
        walk(child)
    return "\n\n".join(parts) + "\n"
